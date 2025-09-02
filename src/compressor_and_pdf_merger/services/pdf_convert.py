from __future__ import annotations
from pathlib import Path
from typing import Optional, Literal
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
import os, math
from .pdf_utils import which, run


def pdf_to_images(
    src_pdf: str | Path,
    out_dir: str | Path,
    *,
    fmt: Literal["jpg","png"]="jpg",
    dpi: int = 144,
    rgb: bool = True,
    page_range: Optional[str] = None,
) -> list[str]:
    out_dir = Path(out_dir); out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(src_pdf))
    pages = _resolve_pages(doc, page_range)
    paths: list[str] = []
    scale = dpi / 72.0
    mat = fitz.Matrix(scale, scale)

    for i, pno in enumerate(pages, 1):
        pix = doc.load_page(pno).get_pixmap(matrix=mat, colorspace=fitz.csRGB if rgb else fitz.csGRAY)
        out = out_dir / f"{Path(src_pdf).stem}_{pno+1:04d}.{fmt}"
        pix.save(str(out))
        paths.append(str(out))
    return paths


def _resolve_pages(doc, rng: Optional[str]) -> list[int]:
    if not rng:
        return list(range(len(doc)))
    out: list[int] = []
    total = len(doc)
    for part in rng.split(","):
        part = part.strip()
        if not part: continue
        if "-" in part:
            a, b = part.split("-", 1)
            start = max(1, int(a)) if a else 1
            end = total if not b else max(1, int(b))
            out.extend([i-1 for i in range(start, end+1)])
        else:
            out.append(max(1, int(part)) - 1)
    return [p for p in out if 0 <= p < total]


def pdf_to_office(
    src_pdf: str | Path,
    out_dir: str | Path,
    *,
    kind: Literal["docx","rtf","xlsx","pptx"]="docx"
) -> str:
    soffice = which("soffice") or which("soffice.exe")
    if not soffice:
        raise RuntimeError("LibreOffice не найден. Поставьте LibreOffice и добавьте soffice в PATH.")
    out_dir = Path(out_dir); out_dir.mkdir(parents=True, exist_ok=True)
    out_name = Path(src_pdf).stem + f".{kind}"
    cmd = [soffice, "--headless", "--convert-to", kind, "--outdir", str(out_dir), str(src_pdf)]
    run(cmd, "LibreOffice")
    out_path = out_dir / out_name
    if not out_path.exists():
        candidates = list(out_dir.glob(f"*.{kind}"))
        if not candidates:
            raise RuntimeError("LibreOffice не смог конвертировать PDF → " + kind.upper())
        out_path = candidates[0]
    return str(out_path)


def pdf_to_pptx_snapshots(
    src_pdf: str | Path,
    out_pptx: str | Path,
    *,
    dpi: int = 144,
    rgb: bool = True,
) -> str:
    doc = fitz.open(str(src_pdf))
    prs = Presentation()
    blank = prs.slide_layouts[6]

    slide_w, slide_h = prs.slide_width, prs.slide_height

    scale = dpi / 72.0
    mat = fitz.Matrix(scale, scale)

    for pno in range(len(doc)):
        page = doc.load_page(pno)
        pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB if rgb else fitz.csGRAY)
        tmp_img = Path(out_pptx).with_name(f"__tmp_slide_{pno:04d}.png")
        pix.save(str(tmp_img))

        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(tmp_img), left=0, top=0, width=slide_w, height=slide_h)
        tmp_img.unlink(missing_ok=True)

    prs.save(str(out_pptx))
    return str(out_pptx)


def pdf_to_text(src_pdf: str | Path, out_txt: str | Path, *, ocr: bool=False, ocr_lang: str="eng") -> str:
    out_txt = Path(out_txt); out_txt.parent.mkdir(parents=True, exist_ok=True)
    if ocr:
        ocrmypdf = which("ocrmypdf")
        if not ocrmypdf:
            raise RuntimeError("Для OCR нужен ocrmypdf (и Tesseract). Установите и добавьте в PATH.")
        tmp_ocr = out_txt.with_suffix(".__ocr__.pdf")
        run([ocrmypdf, "--language", ocr_lang, "--force-ocr", str(src_pdf), str(tmp_ocr)], "ocrmypdf")
        src_pdf = tmp_ocr

    doc = fitz.open(str(src_pdf))
    text = []
    for p in doc:
        text.append(p.get_text("text"))
    out_txt.write_text("\n".join(text), encoding="utf-8")

    if str(src_pdf).endswith(".__ocr__.pdf"):
        Path(src_pdf).unlink(missing_ok=True)
    return str(out_txt)
