from __future__ import annotations
from pathlib import Path
import io
import os
import re
import sys
import shutil
import subprocess
from typing import Optional, Tuple
import mammoth
from xhtml2pdf import pisa
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
import fitz


EMU_PER_INCH = 914_400
NOTO_FILES = {"regular": "NotoSans-Regular.ttf", "bold": "NotoSans-Bold.ttf", "italic": "NotoSans-Italic.ttf", "bolditalic": "NotoSans-BoldItalic.ttf"}
DEJAVU_FILES = {"regular": "DejaVuSans.ttf", "bold": "DejaVuSans-Bold.ttf", "italic": "DejaVuSans-Oblique.ttf", "bolditalic": "DejaVuSans-BoldOblique.ttf"}
_CTRL = "".join(chr(c) for c in range(0, 32) if c not in (9, 10, 13))
_CTRL_RE = re.compile(f"[{re.escape(_CTRL)}]")


def _assets_dir() -> Path:
    if hasattr(sys, "_MEIPASS"):
        return Path(sys._MEIPASS) / "compressor_and_pdf_merger" / "assets"  # type: ignore[attr-defined]
    return Path(__file__).resolve().parents[1] / "assets"


def _fonts_dir() -> Path:
    return _assets_dir() / "fonts"


def _font_path(candidates: Tuple[str, ...]) -> Optional[Path]:
    fonts = _fonts_dir()
    for name in candidates:
        p = fonts / name
        if p.exists():
            return p
    return None


def _register_font_family(family_name: str, regular: Path, bold: Optional[Path] = None, italic: Optional[Path] = None, bolditalic: Optional[Path] = None) -> None:
    pdfmetrics.registerFont(TTFont(family_name, str(regular)))
    if bold:
        pdfmetrics.registerFont(TTFont(f"{family_name}-Bold", str(bold)))
    if italic:
        pdfmetrics.registerFont(TTFont(f"{family_name}-Italic", str(italic)))
    if bolditalic:
        pdfmetrics.registerFont(TTFont(f"{family_name}-BoldItalic", str(bolditalic)))
    pdfmetrics.registerFontFamily(family_name, normal=family_name, bold=f"{family_name}-Bold" if bold else family_name, italic=f"{family_name}-Italic" if italic else family_name, boldItalic=f"{family_name}-BoldItalic" if bolditalic else family_name)


def _ensure_fonts() -> Tuple[str, Optional[Path]]:
    _fonts_dir().mkdir(parents=True, exist_ok=True)
    reg = _font_path((NOTO_FILES["regular"],))
    if reg:
        b = _font_path((NOTO_FILES["bold"],))
        i = _font_path((NOTO_FILES["italic"],))
        bi = _font_path((NOTO_FILES["bolditalic"],))
        _register_font_family("NotoSans", reg, b, i, bi)
        pisa.DEFAULT_FONT = "NotoSans"
        return "NotoSans", reg
    reg = _font_path((DEJAVU_FILES["regular"],))
    if reg:
        b = _font_path((DEJAVU_FILES["bold"],))
        i = _font_path((DEJAVU_FILES["italic"],))
        bi = _font_path((DEJAVU_FILES["bolditalic"],))
        _register_font_family("DejaVuSans", reg, b, i, bi)
        pisa.DEFAULT_FONT = "DejaVuSans"
        return "DejaVuSans", reg
    pisa.DEFAULT_FONT = "Helvetica"
    return "Helvetica", None


def _fallback_css(reportlab_family: str) -> str:
    family_chain = f"{reportlab_family}, 'DejaVu Sans', Arial, sans-serif"
    return "@page{size:A4;margin:18mm}" "html,body{font-size:12pt;}" f"*{{font-family:{family_chain} !important;}}" "h1,h2,h3,h4,h5,h6{font-weight:bold;}" "table{border-collapse:collapse;width:100%;}" "th,td{border:1px solid #ccc;padding:4px;vertical-align:top;text-align:left;}" "thead th{background:#f2f2f2;}" ".pb{page-break-after:always;}" "span.tab,.tab{white-space:pre;}"


def _sanitize_html(s: str) -> str:
    s = _CTRL_RE.sub(" ", s)
    s = s.replace("\t", " ··· ")
    return s


def _find_soffice() -> Optional[str]:
    p = shutil.which("soffice")
    if p:
        return p
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    return None


def _run_soffice_to_pdf(input_path: Path, out_pdf: Path, filter_name: Optional[str] = None, timeout: int = 300) -> bool:
    exe = _find_soffice()
    if not exe:
        return False
    out_dir = out_pdf.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    arg_filter = f"pdf:{filter_name}" if filter_name else "pdf"
    cmd = [exe, "--headless", "--norestore", "--convert-to", arg_filter, "--outdir", str(out_dir), str(input_path)]
    try:
        subprocess.run(cmd, check=True, timeout=timeout, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except Exception:
        return False
    produced = out_dir / (input_path.stem + ".pdf")
    if produced.exists():
        if out_pdf.exists():
            out_pdf.unlink(missing_ok=True)
        produced.replace(out_pdf)
        return True
    return False


def _try_word_docx_to_pdf(docx_path: Path, out_pdf: Path) -> bool:
    try:
        import win32com.client  # type: ignore
        from win32com.client import constants  # type: ignore
    except Exception:
        return False
    word = None
    try:
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(str(docx_path), ReadOnly=True)
        doc.ExportAsFixedFormat(
            OutputFileName=str(out_pdf),
            ExportFormat=17,
            OpenAfterExport=False,
            OptimizeFor=0,
            Range=0,
            Item=0,
            IncludeDocProps=True,
            KeepIRM=True,
            CreateBookmarks=1,
            DocStructureTags=True,
            BitmapMissingFonts=False,
            UseISO19005_1=False,
        )
        doc.Close(False)
        return out_pdf.exists()
    except Exception:
        return False
    finally:
        try:
            if word:
                word.Quit()
        except Exception:
            pass


def _try_excel_xlsx_to_pdf(xlsx_path: Path, out_pdf: Path) -> bool:
    try:
        import win32com.client  # type: ignore
    except Exception:
        return False
    excel = None
    try:
        excel = win32com.client.DispatchEx("Excel.Application")
        excel.Visible = False
        wb = excel.Workbooks.Open(str(xlsx_path))
        wb.ExportAsFixedFormat(0, str(out_pdf))
        wb.Close(SaveChanges=False)
        return out_pdf.exists()
    except Exception:
        return False
    finally:
        try:
            if excel:
                excel.Quit()
        except Exception:
            pass


def _try_pptx_to_pdf_via_powerpoint(pptx_path: Path, out_pdf: Path) -> bool:
    try:
        import win32com.client  # type: ignore
    except Exception:
        return False
    app = None
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = True
        pres = app.Presentations.Open(str(pptx_path), WithWindow=False)
        pres.ExportAsFixedFormat(str(out_pdf), 2, 2)
        pres.Close()
        return out_pdf.exists()
    except Exception:
        return False
    finally:
        try:
            if app:
                app.Quit()
        except Exception:
            pass


def docx_to_pdf_basic(docx_path: str | Path, out_pdf: str | Path) -> str:
    src = Path(docx_path)
    dst = Path(out_pdf)
    dst.parent.mkdir(parents=True, exist_ok=True)
    if _try_word_docx_to_pdf(src, dst):
        return str(dst)
    if _run_soffice_to_pdf(src, dst, "writer_pdf_Export"):
        return str(dst)
    family, _ = _ensure_fonts()
    css = _fallback_css(family)
    with open(src, "rb") as f:
        try:
            conv = mammoth.images.inline(mammoth.images.img_element(src=mammoth.images.data_uri))
            res = mammoth.convert_to_html(f, convert_image=conv)
        except Exception:
            res = mammoth.convert_to_html(f)
    html = f"<html><head><meta charset='utf-8'><style>{css}</style></head><body>{res.value}</body></html>"
    html = _sanitize_html(html)
    with open(dst, "wb") as out:
        pisa.CreatePDF(html, dest=out)
    return str(dst)


def xlsx_to_pdf_basic(xlsx_path: str | Path, out_pdf: str | Path) -> str:
    src = Path(xlsx_path)
    dst = Path(out_pdf)
    dst.parent.mkdir(parents=True, exist_ok=True)
    if _try_excel_xlsx_to_pdf(src, dst):
        return str(dst)
    if _run_soffice_to_pdf(src, dst, "calc_pdf_Export"):
        return str(dst)
    family, _ = _ensure_fonts()
    css = _fallback_css(family)
    xls = pd.ExcelFile(src)
    parts = [f"<html><head><meta charset='utf-8'><style>{css}</style></head><body>"]
    for i, sheet in enumerate(xls.sheet_names):
        df = xls.parse(sheet, dtype=str).fillna("")
        parts.append(f"<h2>{sheet}</h2>")
        parts.append(df.to_html(index=False, border=0))
        if i < len(xls.sheet_names) - 1:
            parts.append("<div class='pb'></div>")
    parts.append("</body></html>")
    html = _sanitize_html("".join(parts))
    with open(dst, "wb") as out:
        pisa.CreatePDF(html, dest=out)
    return str(dst)


def _emu_to_px(val_emu: int, dpi: int) -> int:
    return max(1, int((val_emu / EMU_PER_INCH) * dpi))


def _pt_to_px(pt: float, dpi: int) -> int:
    return max(8, int(round(pt / 72.0 * dpi)))


def _shape_box(shape, dpi: int, offset: Tuple[int, int] = (0, 0)) -> Optional[Tuple[int, int, int, int]]:
    try:
        l = _emu_to_px(int(shape.left), dpi) + offset[0]
        t = _emu_to_px(int(shape.top), dpi) + offset[1]
        w = _emu_to_px(int(shape.width), dpi)
        h = _emu_to_px(int(shape.height), dpi)
        return l, t, w, h
    except Exception:
        return None


def _render_textframe(draw, shape, dpi: int, pil_font_regular: ImageFont.FreeTypeFont, offset=(0, 0)) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    box = _shape_box(shape, dpi, offset)
    if not box:
        return
    l, t, w, h = box
    y = t
    try:
        for p in shape.text_frame.paragraphs:
            runs = getattr(p, "runs", []) or []
            txt = "".join([getattr(r, "text", "") for r in runs]) or (p.text or "")
            if not txt.strip():
                continue
            sz_pt = next((float(r.font.size.pt) for r in runs if getattr(r.font, "size", None)), 18.0)
            try:
                font = ImageFont.truetype(pil_font_regular.path, _pt_to_px(sz_pt, dpi))  # type: ignore[attr-defined]
            except Exception:
                font = pil_font_regular
            draw.multiline_text((l + 6, y + 4), txt, fill=(0, 0, 0), font=font)
            y += getattr(font, "size", 14) + 6
    except Exception:
        try:
            draw.multiline_text((l + 6, y + 4), shape.text_frame.text or "", fill=(0, 0, 0), font=pil_font_regular)
        except Exception:
            pass


def _render_picture(canvas: Image.Image, shape, dpi: int, offset=(0, 0)) -> None:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE or not hasattr(shape, "image"):
        return
    box = _shape_box(shape, dpi, offset)
    if not box:
        return
    l, t, w, h = box
    try:
        pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
        pic = pic.resize((max(1, w), max(1, h)))
        canvas.paste(pic, (l, t))
    except Exception:
        pass


def _render_autoshape(draw, shape, dpi: int, offset=(0, 0)) -> None:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return
    box = _shape_box(shape, dpi, offset)
    if not box:
        return
    l, t, w, h = box
    try:
        color = (230, 230, 230)
        fc = getattr(shape, "fill", None)
        if fc and getattr(fc, "type", None) is not None and hasattr(fc, "fore_color"):
            col = getattr(fc.fore_color, "rgb", None)
            if col is not None:
                color = (int(col[0]), int(col[1]), int(col[2]))
        draw.rectangle([l, t, l + w, t + h], fill=color, outline=(180, 180, 180))
    except Exception:
        draw.rectangle([l, t, l + w, t + h], outline=(180, 180, 180))


def _render_table(draw, shape, dpi: int, pil_font_regular: ImageFont.FreeTypeFont, offset=(0, 0)) -> None:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.TABLE:
        return
    box = _shape_box(shape, dpi, offset)
    if not box:
        return
    l, t, w, h = box
    try:
        rows = shape.table.rows
        cols = shape.table.columns
        rh = max(1, h // max(1, len(rows)))
        cw = max(1, w // max(1, len(cols)))
        for r, _ in enumerate(rows):
            for c, _ in enumerate(cols):
                x0 = l + c * cw
                y0 = t + r * rh
                x1 = x0 + cw
                y1 = y0 + rh
                draw.rectangle([x0, y0, x1, y1], outline=(180, 180, 180))
                try:
                    cell = shape.table.cell(r, c)
                    txt = (cell.text or "").strip()
                    if txt:
                        draw.multiline_text((x0 + 4, y0 + 4), txt, fill=(0, 0, 0), font=pil_font_regular)
                except Exception:
                    pass
    except Exception:
        pass


def _render_group(canvas, draw, shape, dpi: int, pil_font_regular: ImageFont.FreeTypeFont, offset=(0, 0)) -> None:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.GROUP:
        return
    gbox = _shape_box(shape, dpi, offset)
    base_off = (gbox[0], gbox[1]) if gbox else offset
    try:
        for child in shape.shapes:
            _render_one(canvas, draw, child, dpi, pil_font_regular, base_off)
    except Exception:
        pass


def _render_one(canvas, draw, shape, dpi: int, pil_font_regular: ImageFont.FreeTypeFont, offset=(0, 0)) -> None:
    try:
        st = getattr(shape, "shape_type", None)
        if st == MSO_SHAPE_TYPE.GROUP:
            _render_group(canvas, draw, shape, dpi, pil_font_regular, offset)
            return
        if st == MSO_SHAPE_TYPE.PICTURE:
            _render_picture(canvas, shape, dpi, offset)
        elif st == MSO_SHAPE_TYPE.TABLE:
            _render_table(draw, shape, dpi, pil_font_regular, offset)
        elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            _render_autoshape(draw, shape, dpi, offset)
        if getattr(shape, "has_text_frame", False):
            _render_textframe(draw, shape, dpi, pil_font_regular, offset)
    except Exception:
        pass


def pptx_to_pdf_basic(pptx_path: str | Path, out_pdf: str | Path, dpi: int = 150) -> str:
    src = Path(pptx_path)
    dst = Path(out_pdf)
    dst.parent.mkdir(parents=True, exist_ok=True)
    if _try_pptx_to_pdf_via_powerpoint(src, dst):
        return str(dst)
    if _run_soffice_to_pdf(src, dst, "impress_pdf_Export"):
        return str(dst)
    _, pil_ttf = _ensure_fonts()
    if pil_ttf and Path(pil_ttf).exists():
        pil_font = ImageFont.truetype(str(pil_ttf), size=max(12, dpi // 11))
    else:
        try:
            pil_font = ImageFont.truetype("DejaVuSans.ttf", size=max(12, dpi // 11))
        except Exception:
            pil_font = ImageFont.load_default()
    prs = Presentation(str(src))
    page_w_px = _emu_to_px(int(prs.slide_width), dpi)
    page_h_px = _emu_to_px(int(prs.slide_height), dpi)
    doc = fitz.open()
    for slide in prs.slides:
        img = Image.new("RGB", (page_w_px, page_h_px), "white")
        draw = ImageDraw.Draw(img)
        try:
            for shape in slide.shapes:
                _render_one(img, draw, shape, dpi, pil_font, (0, 0))
        except Exception:
            pass
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        page = doc.new_page(width=page_w_px, height=page_h_px)
        rect = fitz.Rect(0, 0, page_w_px, page_h_px)
        page.insert_image(rect, stream=buf.getvalue())
    doc.save(str(dst))
    doc.close()
    return str(dst)


__all__ = ["docx_to_pdf_basic", "xlsx_to_pdf_basic", "pptx_to_pdf_basic"]
